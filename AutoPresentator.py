from tkinter import *
from tkinter import messagebox
import pyttsx3
import time
import subprocess
from pptx import Presentation

import glob
from pptx.enum.action import PP_ACTION
from tkinter import filedialog
from ttkthemes import themed_tk as tk
import threading
import pyautogui as pg
import os
# configuring the audio and tts engine
engine=pyttsx3.init()
voices=engine.getProperty('voices')
engine.setProperty('voice',voices[1])


def entered(event):
    nButton.config(bg="#343434")
    nButton.config(fg="#ffffff")
def left(event):
    nButton.config(fg="#343434")
    nButton.config(bg="#ffffff")
def speak(txtt):    #a function that takes txtt and speaks it
    engine.say(txtt)
    engine.runAndWait()

subprocess.call(["open","-a",'Microsoft PowerPoint' ])
# GUI stuffs tkinter and themes
root=tk.ThemedTk()
root.get_themes()
root.set_theme('plastik')
root.title("Bishesh-Auto_Presentator")
root.iconbitmap('favicon.ico')

# defining some string varaiable in gui
texts=StringVar()
textn=StringVar()
# fetching some data from the gui text entry and the main function

def open_file():
    global a
    root.a=filedialog.askopenfilename(initialdir="/",title="Open Presentation",filetypes= [('Presentation Files', '*.pptx'),('Presentation Files','*ppt')] )
    a=root.a

    nlayout=Label(root,text=a,font="Helvetica 18 bold italic").pack()
       
    

nopfile=Button(root,text="Open File",command=open_file).pack()

def click():
    
    time.sleep(5)
    try:

        for eachfile in glob.glob(a):
               
                t1=threading.Thread(target=subprocess.call(["open",a]))
                t1.start()
                #pg.hotkey('command','shift','enter')
                time.sleep(10)

                prs = Presentation(eachfile)
            
                
                

                i=1
                

                
                for slide in prs.slides:
                    slds=len(list(prs.slides))
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            b=len(shape.text)
                            speak(shape.text)
                    # speak("You may explain")
                    
                    # time.sleep(int(texts.get()))
                    speak("Next SLide ")
                    if i>=slds:
                        pg.hotkey('right')
                        pg.hotkey("esc","command",".")
                        
                    else:
                        pg.hotkey("right")
                        
                        i+=1

                    

    except Exception as e:
        print("Error")
        with open("a.txt","w") as f:
            f.write(str(e))
    finally:
        root.destroy()
        

                
    

timefex=Entry(root,textvariable=texts,width=5).pack()


# nEntry=Entry(root,textvariable=textx,width=60).pack()
nButton=Button(root,command=click,text="Start",bg="#ffffff",fg="#343434",font="Helvetica 18 bold italic")
nButton.pack()
nButton.bind("<Enter>",entered)
nButton.bind("<Leave>",left)

root.layou=Label(text="open a .pptx file, start the presentation and wait..",font="Helvetica  18 bold italic").pack()

root.layoutn=Label(text="Made BY BISHESH BOHORA @2020 V1.1",font="Helvetica  18 bold italic").pack()



root.mainloop()